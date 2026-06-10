"""
Slide API — parse all presentation formats + AI generation via ppt-generator skill
Run from project root:
  cd server && pip install -r requirements.txt && uvicorn app:app --host 0.0.0.0 --port 8000

Env:
  OPENAI_API_KEY      — required for /api/generate
  OPENAI_BASE_URL     — optional, default OpenAI
  OPENAI_MODEL        — optional, default gpt-4o-mini
  LibreOffice         — required for .ppt/.pps/.odp parse fallback
"""
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv

load_dotenv(Path(__file__).resolve().parent / ".env")

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pptx import Presentation

from ai_generator import ai_configured, generate_presentation_html, get_model_presets

ALLOWED = {".ppt", ".pptx", ".pps", ".ppsx", ".odp"}
SOFFICE = None
ROOT = Path(__file__).resolve().parent.parent


def find_soffice() -> Optional[str]:
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    for path in (
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ):
        if path.exists():
            return str(path)
    return None


SOFFICE = find_soffice()

app = FastAPI(title="Slide API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


def convert_to_pptx(src: Path, out_dir: Path) -> Path:
    if not SOFFICE:
        raise HTTPException(
            503,
            "未安装 LibreOffice，无法解析 .ppt/.pps/.odp。请安装 LibreOffice 或上传 .pptx",
        )
    cmd = [SOFFICE, "--headless", "--convert-to", "pptx", "--outdir", str(out_dir), str(src)]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    except subprocess.TimeoutExpired as e:
        raise HTTPException(504, "文件转换超时") from e
    if result.returncode != 0:
        raise HTTPException(422, f"LibreOffice 转换失败: {result.stderr[:200]}")

    pptx = out_dir / f"{src.stem}.pptx"
    if not pptx.exists():
        candidates = list(out_dir.glob("*.pptx"))
        if not candidates:
            raise HTTPException(422, "转换后未找到 pptx")
        pptx = candidates[0]
    return pptx


def shape_texts(shape) -> list[str]:
    texts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = (para.text or "").strip()
            if t:
                texts.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                t = (cell.text or "").strip()
                if t:
                    texts.append(t)
    return texts


def extract_slides_from_pptx(pptx_path: Path) -> list[dict]:
    prs = Presentation(str(pptx_path))
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            texts.extend(shape_texts(shape))
        content = "\n".join(texts)
        slides.append(
            {
                "index": len(slides) + 1,
                "texts": texts,
                "content": content,
                "preview": content.replace("\n", " · ")[:80] or "(空白页)",
                "title": texts[0] if texts else "",
                "body": texts[1:] if len(texts) > 1 else [],
            }
        )
    return slides


async def parse_file_bytes(raw: bytes, filename: str) -> list[dict]:
    ext = Path(filename or "upload.pptx").suffix.lower()
    if ext not in ALLOWED:
        raise HTTPException(400, f"不支持 {ext}，支持: {', '.join(sorted(ALLOWED))}")
    if len(raw) > 50 * 1024 * 1024:
        raise HTTPException(413, "文件不能超过 50MB")
    if not raw:
        raise HTTPException(400, "空文件")

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        src = tmp_path / Path(filename or "upload.pptx").name
        src.write_bytes(raw)

        if ext in (".pptx", ".ppsx"):
            target = src
        else:
            target = convert_to_pptx(src, tmp_path)

        slides = extract_slides_from_pptx(target)
        if not slides:
            raise HTTPException(422, "未能提取到文字，请确认幻灯片含可编辑文本")
        return slides


@app.get("/api/health")
def health():
    return {
        "ok": True,
        "server_default_ai": ai_configured(),
        "libreoffice": SOFFICE is not None,
        "formats": sorted(ALLOWED),
    }


@app.get("/api/models")
def list_models():
    return {"presets": get_model_presets()}


@app.post("/api/parse")
async def parse_upload(file: UploadFile = File(...)):
    raw = await file.read()
    slides = await parse_file_bytes(raw, file.filename or "upload.pptx")
    return {"slides": slides}


@app.post("/api/generate")
async def generate_upload(
    file: UploadFile = File(...),
    api_key: str = Form(""),
    base_url: str = Form(""),
    model: str = Form(""),
):
    """Upload PPT → extract script → user-selected LLM + ppt-generator skill → HTML"""
    raw = await file.read()
    name = file.filename or "presentation.pptx"
    slides = await parse_file_bytes(raw, name)
    html = await generate_presentation_html(
        slides,
        name,
        api_key=api_key.strip() or None,
        base_url=base_url.strip() or None,
        model=model.strip() or None,
    )
    return {
        "html": html,
        "slides_count": len(slides),
        "file_name": name,
        "engine": "ai+skill",
        "model": model.strip() or None,
    }


if ROOT.joinpath("index.html").exists():
    app.mount("/", StaticFiles(directory=str(ROOT), html=True), name="site")
